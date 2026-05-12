import os
import subprocess
from pdf2docx import Converter
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt as PptxPt
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN

def _require_file(path: str, name: str):
    if not path or not os.path.exists(path):
        raise RuntimeError(f"{name} topilmadi: {path}")

def pdf_to_docx(pdf_path: str, out_docx: str):
    if not os.path.exists(pdf_path):
        raise RuntimeError(f"PDF topilmadi: {pdf_path}")
    cv = Converter(pdf_path)
    try:
        cv.convert(out_docx, start=0, end=None)
    finally:
        cv.close()

def docx_to_pdf(libreoffice_path: str, docx_path: str, out_dir: str) -> str:
    _require_file(libreoffice_path, "LibreOffice (soffice.exe)")
    if not os.path.exists(docx_path):
        raise RuntimeError(f"DOCX topilmadi: {docx_path}")
    os.makedirs(out_dir, exist_ok=True)

    cmd = [libreoffice_path, "--headless", "--convert-to", "pdf", "--outdir", out_dir, docx_path]
    subprocess.check_call(cmd)

    base = os.path.splitext(os.path.basename(docx_path))[0]
    pdf_path = os.path.join(out_dir, base + ".pdf")
    if not os.path.exists(pdf_path):
        raise RuntimeError("LibreOffice PDF chiqarmadi.")
    return pdf_path

def text_to_docx(text: str, out_docx: str, title: str | None = None):
    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for line in (text or "").splitlines():
        clean_line = line.strip()
        if not clean_line: continue
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = p.add_run(clean_line)
        run.font.name = 'Times New Roman'
        run.font.size = DocxPt(14)
        # Word dasturi shriftni aniq tanishi uchun XML sozlamasi
        run._element.rPr.rFonts.set(qn('w:ascii'), 'Times New Roman')
        run._element.rPr.rFonts.set(qn('w:hAnsi'), 'Times New Roman')
    doc.save(out_docx)

def text_to_pptx(text: str, out_pptx: str, title: str = "Generated Slides"):
    prs = Presentation()
    # Matnni barcha qatorlarga bo'lamiz
    lines = (text or "").splitlines()
    if not lines:
        lines = [" "]

    blank_slide_layout = prs.slide_layouts[6] # Odatda butunlay bo'sh slayd
    
    # 20pt shriftda va 0.5 dyuymli chekkalarda slaydga taxminan 11 qator sig'adi
    MAX_LINES = 11 
    current_slide = None
    current_text_frame = None
    line_count = 0

    for line in lines:
        text_line = line.strip()
        if not text_line:
            # Bo'sh qator tashlash
            if current_text_frame:
                current_text_frame.add_paragraph()
                line_count += 1
            continue

        # Aqlli wrap hisobi: 20pt Times New Roman uchun bir qatorga taxminan 65-70 belgi sig'adi
        needed = max(1, (len(text_line) // 65) + 1)

        # Yangi slayd ochish kerakligini tekshiramiz
        if current_slide is None or (line_count + needed > MAX_LINES):
            current_slide = prs.slides.add_slide(blank_slide_layout)
            
            # Slayd o'lchami odatda 10x7.5 dyuym. Matn qutisini chekkalar bilan joylaymiz.
            left, top = PptxInches(0.5), PptxInches(0.5)
            width, height = PptxInches(9.0), PptxInches(6.5)
            txBox = current_slide.shapes.add_textbox(left, top, width, height)
            current_text_frame = txBox.text_frame
            current_text_frame.word_wrap = True
            line_count = 0

        # Paragraf qo'shish va formatlash
        p = current_text_frame.add_paragraph()
        p.alignment = PP_ALIGN.JUSTIFY
        run = p.add_run()
        run.text = text_line
        run.font.name = 'Times New Roman'
        run.font.size = PptxPt(20)
        line_count += needed

    prs.save(out_pptx)

def images_to_docx_embed(image_paths: list[str], out_docx: str):
    if not image_paths:
        raise RuntimeError("Rasmlar yo'q.")
    doc = Document()
    
    # Sahifa chekkalarini (margins) 0.5 dyuymga o'rnatamiz
    section = doc.sections[0]
    section.top_margin = DocxInches(0.5)
    section.bottom_margin = DocxInches(0.5)
    section.left_margin = DocxInches(0.5)
    section.right_margin = DocxInches(0.5)

    for img_p in image_paths:
        if os.path.exists(img_p):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            # Rasmni sahifa kengligiga moslab (6.5 dyuym) joylaymiz
            run.add_picture(img_p, width=DocxInches(6.5))
    doc.save(out_docx)
